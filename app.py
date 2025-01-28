import os
import gdown
from google.colab import drive
import openai
import nltk
from pptx import Presentation
import openpyxl
import fitz  # PyMuPDF
import numpy as np

# Download required NLTK data for processing language:
# punkt -> sentence splitting, averaged_perceptron_tagger -> pos, stopwords
nltk.download('punkt')
nltk.download('averaged_perceptron_tagger')
nltk.download('stopwords')

# Step 1: Mount Google Drive
drive.mount('/content/drive')

# Set base path to the Google Drive folder where "Upload" is stored
base_path = '/content/drive/MyDrive/Upload'  # Adjust path if necessary
pdf_folder = os.path.join(base_path, "pdfs")
ppt_folder = os.path.join(base_path, "ppts")
excel_folder = os.path.join(base_path, "excels")

# Authenticate OpenAI API
openai.api_key = os.getenv("OPENAI_API_KEY")

# Function to download files from Google Drive
def download_files_from_drive(folder_url, save_path):
    # Extracting file IDs manually or from the folder URL
    # Example of direct file links (you can generate these from your folder)
    file_urls = [
        "https://drive.google.com/uc?id=1BfgdaqBgQL-ku0l-iCQTYA4i5WKQw32V"  # Replace with actual file IDs
        # Add more file IDs here
    ]
    
    for file_url in file_urls:
        save_file_path = os.path.join(save_path, file_url.split('=')[-1] + ".pdf")  # Or adjust extension as needed
        gdown.download(file_url, save_file_path, quiet=False)
        print(f"Downloaded {save_file_path}")

# Download files from the folder
download_files_from_drive(pdf_folder, save_path="./downloads")

# Step 2: Text Extraction Functions
# Extract text from PDF
def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

# Extract text from PPT
def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Extract text from Excel
def extract_text_from_excel(excel_path):
    wb = openpyxl.load_workbook(excel_path)
    sheet = wb.active
    text = ""
    for row in sheet.iter_rows():
        for cell in row:
            text += str(cell.value) + "\n"
    return text

# Step 3: Text Preprocessing and Chunking
# Clean the text
def clean_text(text):
    text = text.replace("\x00", "")
    text = text.encode("ascii", "ignore").decode()
    text = " ".join(text.split())
    return text

# Tokenizes text into sentences, if current chunk has space for new sentence add, else new chunk
# Keeps overlap sentences to keep context
def split_text_into_chunks(text, max_tokens=500, overlap=100):
    sentences = nltk.tokenize.sent_tokenize(text)
    chunks = []
    current_chunk = []
    current_length = 0

    for sentence in sentences:
        sentence_length = len(sentence.split())
        if current_length + sentence_length > max_tokens:
            # Create chunk with overlap
            chunk_text = ' '.join(current_chunk)
            chunks.append(chunk_text)

            # Keep last few sentences for overlap
            overlap_words = chunk_text.split()[-overlap:]
            current_chunk = [' '.join(overlap_words), sentence]
            current_length = len(overlap_words) + sentence_length
        else:
            current_chunk.append(sentence)
            current_length += sentence_length

    if current_chunk:
        chunks.append(' '.join(current_chunk))

    return chunks

# Send small batches of text to OpenAI to embed
def generate_embeddings(texts, batch_size=10):
    embeddings = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        try:
            response = openai.Embedding.create(
                model="text-embedding-ada-002",
                input=batch
            )
            embeddings.extend([embedding["embedding"] for embedding in response["data"]])
        except Exception as e:
            print(f"Error generating embeddings for batch {i}-{i+batch_size}: {e}")
    return embeddings

# Store document chunks and their embeddings
document_store = {}

# Process and upload from folder
def process_and_upload_from_folder(folder_path, doc_type, chunking_strategy):
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        if filename.endswith(".pdf") and doc_type == "pdf":
            text = extract_text_from_pdf(file_path)
        elif filename.endswith(".pptx") and doc_type == "ppt":
            text = extract_text_from_ppt(file_path)
        elif filename.endswith(".xlsx") and doc_type == "excel":
            text = extract_text_from_excel(file_path)
        else:
            continue

        # Clean and chunk the text
        clean_doc_text = clean_text(text)
        text_chunks = chunking_strategy(clean_doc_text)
        embeddings = generate_embeddings(text_chunks)

        document_store[filename] = {
            "chunks": text_chunks[:len(embeddings)],
            "embeddings": embeddings,
            "source": doc_type
        }
        print(f"Uploaded {filename} ({doc_type.upper()}) to the memory store.")

# Cosine similarity
def cosine_similarity(vec1, vec2):
    vec1_norm = np.linalg.norm(vec1)
    vec2_norm = np.linalg.norm(vec2)

    if vec1_norm == 0 or vec2_norm == 0:
        return 0.0

    return np.dot(vec1, vec2) / (vec1_norm * vec2_norm)

# Retrieve relevant chunks based on query
def retrieve_relevant_chunks(query, top_k=3):
    query_embedding = generate_embeddings([query])[0]
    similarities = []

    for doc_name, doc_data in document_store.items():
        for chunk, chunk_embedding in zip(doc_data["chunks"], doc_data["embeddings"]):
            similarity = cosine_similarity(query_embedding, chunk_embedding)
            similarities.append((chunk, similarity, doc_name))

    relevant_chunks = sorted(similarities, key=lambda x: x[1], reverse=True)[:top_k]
    return [(chunk, doc_name) for chunk, _, doc_name in relevant_chunks]

# Chat with the assistant
def chat_with_assistant(query):
    relevant_chunks = retrieve_relevant_chunks(query)
    print(f"Relevant chunks for query '{query}':")
    for chunk, doc in relevant_chunks:
        print(f"Source ({doc}): {chunk}")
    context = "\n\n".join([f"Source ({doc}): {chunk}" for chunk, doc in relevant_chunks])

    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[{
            "role": "system",
            "content": """You are a precise assistant that answers questions based strictly on the provided context.
            Rules:
            1. Use ONLY information from the context
            2. Keep exact terminology and steps from the source
            3. If multiple sources have different information, specify which source you're using
            4. If information isn't in the context, say "I don't have enough information"
            5. For procedures, list exact steps in order
            6. Include specific buttons, links, and UI elements mentioned in the source"""
        }, {
            "role": "user",
            "content": f"Context:\n{context}\n\nQuestion: {query}"
        }],
        temperature=0,
        max_tokens=500
    )

    return response.choices[0].message.content.strip()

# Process documents
def process_and_upload_all():
    process_and_upload_from_folder(pdf_folder, "pdf", split_text_into_chunks)
    process_and_upload_from_folder(ppt_folder, "ppt", split_text_into_chunks)
    process_and_upload_from_folder(excel_folder, "excel", split_text_into_chunks)

# Run the processing
process_and_upload_all()

# Test the improved system
query = "what is the role of ministry of health in campaign digitization"
answer = chat_with_assistant(query)
print(f"Assistant's answer: {answer}")
